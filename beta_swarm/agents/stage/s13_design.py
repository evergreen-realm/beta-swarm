"""
S13 — Design Agent (full rewrite).
Generates presentation materials, branding guide, feature map, and SVG infographic.
Uses LLM + guaranteed fallbacks. Syncs to SQLiteBrain.
"""
from beta_swarm.agents.base import BaseAgent
import json, os, re, logging
from typing import Dict, Any, List

logger = logging.getLogger(__name__)


class S13DesignAgent(BaseAgent):
    def __init__(self, brain=None):
        super().__init__("s13_design", "Design Agent", "Stage 13: Visual Design", brain)

    def _get_default_next_stage(self):
        return "x1_code_review"

    def execute(self, task: Dict[str, Any]) -> Dict[str, Any]:
        project_id   = task.get("project_id", "default")
        project_path = task.get("project_path", f"./projects/{project_id}")
        s3_out       = task.get("s3_prd", {})
        prd          = s3_out.get("prd") or task.get("prd") or {}
        title        = prd.get("metadata", {}).get("title", "App")
        features     = prd.get("functional_requirements", [])
        preview_url  = task.get("preview_url", "http://localhost:8000")

        self._log_handover(f"S13 started. title={title}")

        design_dir = os.path.join(project_path, "design")
        os.makedirs(design_dir, exist_ok=True)
        files_written: List[str] = []

        def write(rel: str, content: str):
            p = os.path.join(project_path, rel)
            os.makedirs(os.path.dirname(p), exist_ok=True)
            with open(p, "w", encoding="utf-8") as f:
                f.write(content)
            if rel not in files_written:
                files_written.append(rel)

        # ── Presentation (PRESENTATION.md) ─────────────────────────── #
        pres_prompt = f"""Create a stakeholder presentation outline for "{title}".
Include 6 slides: Title, Problem, Solution, Architecture, Demo, Roadmap.
For each slide include: title, key points, speaker notes.
Output markdown."""
        pres = self._call_llm(pres_prompt, task_type="s13_design")
        m = re.search(r'```(?:markdown|md)?\s*\n?(.*?)\n?```', pres, re.DOTALL)
        if m: pres = m.group(1).strip()
        if len(pres) < 100: pres = self._fallback_presentation(title, features, preview_url)
        write("design/PRESENTATION.md", pres)

        # ── Branding Guide (BRANDING.md) ───────────────────────────── #
        brand_prompt = f"""Create a branding guide for "{title}".
Include: primary/secondary colors (hex), typography (fonts + weights), logo concept, tone of voice.
Design theme: dark mode glassmorphism with neon cyan accents.
Output markdown."""
        brand = self._call_llm(brand_prompt, task_type="s13_design")
        m = re.search(r'```(?:markdown|md)?\s*\n?(.*?)\n?```', brand, re.DOTALL)
        if m: brand = m.group(1).strip()
        if len(brand) < 100: brand = self._fallback_branding(title)
        write("design/BRANDING.md", brand)

        # ── Feature Map (FEATURE_MAP.md) ───────────────────────────── #
        feat_prompt = f"""Create a feature map for "{title}".
Features: {json.dumps(features[:8])}
Organize into: Core Features, Nice-to-Have, Future Roadmap.
Use a markdown table or tree structure."""
        feat = self._call_llm(feat_prompt, task_type="s13_design")
        m = re.search(r'```(?:markdown|md)?\s*\n?(.*?)\n?```', feat, re.DOTALL)
        if m: feat = m.group(1).strip()
        if len(feat) < 100: feat = self._fallback_feature_map(title, features)
        write("design/FEATURE_MAP.md", feat)

        # ── SVG Infographic ────────────────────────────────────────── #
        svg_prompt = f"""Generate a complete SVG infographic showing the Beta Swarm pipeline for "{title}".
13 stages as horizontal flow: S1→S2→...→S13 with icons and labels.
Use: dark background #050a0f, neon cyan lines #00f2ff, white labels, rounded boxes.
Output ONLY the raw SVG XML (no markdown wrapping)."""
        svg_raw = self._call_llm(svg_prompt, task_type="s13_design")
        m = re.search(r'```(?:xml|svg)?\s*\n?(<svg.*?</svg>)', svg_raw, re.DOTALL | re.IGNORECASE)
        if m: svg_raw = m.group(1).strip()
        if not svg_raw.strip().startswith("<svg"):
            svg_raw = self._fallback_svg(title)
        write("design/pipeline_infographic.svg", svg_raw)

        # ── PPTX (python-pptx if available) ───────────────────────── #
        pptx_path = os.path.join(project_path, "design", "deck.pptx")
        try:
            from pptx import Presentation
            prs = Presentation()
            slides_data = [
                (title, "Beta Swarm — Autonomous AI Development"),
                ("The Problem", "Software development is slow and fragmented"),
                ("Our Solution", f"{title}: AI-powered end-to-end automation"),
                ("Architecture", "13-stage autonomous pipeline"),
                ("Live Demo", f"Preview: {preview_url}"),
                ("Roadmap", "Phase 1: Core → Phase 2: Self-Growing Brain"),
            ]
            for slide_title, subtitle in slides_data:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_title
                slide.placeholders[1].text = subtitle
            prs.save(pptx_path)
            files_written.append("design/deck.pptx")
            logger.info("[S13] PPTX deck generated.")
        except ImportError:
            logger.warning("[S13] python-pptx not installed — skipping PPTX (PRESENTATION.md generated instead)")
        except Exception as e:
            logger.warning(f"[S13] PPTX generation failed (non-fatal): {e}")

        # ── SQLiteBrain sync ─────────────────────────────────────────#
        try:
            from beta_swarm.brain.sqlite_brain import SQLiteBrain
            db = SQLiteBrain.get_instance()
            db.register_agent("s13_design", "Design Agent", "Stage 13")
            db.store_artifact(agent_id="s13_design", project=title, stage="S13",
                              data=f"Design assets: {', '.join(files_written)}")
        except Exception as e:
            logger.warning(f"[S13] SQLiteBrain sync (non-fatal): {e}")

        artifact = {"design_files": files_written}
        artifact_path = f"./projects/{project_id}/s13_design_output.json"
        os.makedirs(os.path.dirname(artifact_path), exist_ok=True)
        with open(artifact_path, "w", encoding="utf-8") as f:
            json.dump(artifact, f, indent=2)

        self._log_handover(f"S13 completed. {len(files_written)} design files.")

        return {
            "status": "complete",
            "design_files": files_written,
            "artifact": artifact,
            "artifact_path": artifact_path,
            "next_stage": task.get("next_stage") or self._get_default_next_stage()
        }

    # ── Fallbacks ──────────────────────────────────────────────────────── #
    def _fallback_presentation(self, title: str, features: list, preview_url: str) -> str:
        feat_lines = "\n".join([f"  - {f}" for f in features[:5]])
        feat_lines_val = feat_lines or "  - Core functionality\n  - Modern UI\n  - Secure API"
        return f"""# {title} — Stakeholder Presentation

## Slide 1: Title
**{title}**
*Beta Swarm — Autonomous AI Development Platform*

## Slide 2: The Problem
- Software development is slow and expensive
- Teams spend 70% of time on boilerplate
- Inconsistent quality across projects

## Slide 3: Our Solution
- Autonomous 13-stage AI pipeline
- From idea to deployed app in minutes
- Self-growing brain that improves over time

## Slide 4: Key Features
{feat_lines_val}

## Slide 5: Live Demo
- Preview URL: {preview_url}
- Health check: {preview_url}/health

## Slide 6: Roadmap
- **Phase 1** (Now): Core pipeline complete
- **Phase 2**: Self-growing brain
- **Phase 3**: Multi-agent collaboration
"""

    def _fallback_branding(self, title: str) -> str:
        return f"""# {title} — Branding Guide

## Color Palette
| Role       | Hex       | Usage                    |
|------------|-----------|--------------------------|
| Background | `#050a0f` | Page background          |
| Primary    | `#00f2ff` | Buttons, links, accents  |
| Secondary  | `#6366f1` | Cards, hover states      |
| Text       | `#e2e8f0` | Body text                |
| Muted      | `#64748b` | Secondary text, borders  |

## Typography
- **Headings**: Inter Bold (700) — [Google Fonts](https://fonts.google.com/specimen/Inter)
- **Body**: Inter Regular (400)
- **Code**: JetBrains Mono (400)

## Design Principles
1. **Dark by default** — dark mode first
2. **Glassmorphism** — frosted glass cards (`backdrop-filter: blur`)
3. **Neon accents** — cyan glow on interactive elements
4. **Smooth transitions** — 200ms ease on hover states

## Logo Concept
- Icon: hexagonal neural network pattern
- Color: gradient from `#00f2ff` to `#6366f1`
- Font: Inter ExtraBold
"""

    def _fallback_feature_map(self, title: str, features: list) -> str:
        core = features[:5] if features else ["Core functionality", "User authentication", "CRUD operations"]
        nice = ["Advanced search", "Real-time updates", "Export to PDF"]
        future = ["Mobile app", "AI recommendations", "Third-party integrations"]
        core_lines  = "\n".join([f"  - ✅ {f}" for f in core])
        nice_lines  = "\n".join([f"  - 🔶 {f}" for f in nice])
        future_lines = "\n".join([f"  - 🔮 {f}" for f in future])
        return f"""# {title} — Feature Map

## Core Features (MVP)
{core_lines}

## Nice-to-Have (v1.1)
{nice_lines}

## Future Roadmap (v2.0+)
{future_lines}

## Feature Priority Matrix
| Feature           | Impact | Effort | Priority |
|-------------------|--------|--------|----------|
| Core CRUD         | High   | Low    | P0       |
| Auth system       | High   | Medium | P0       |
| Search            | Medium | Medium | P1       |
| Real-time updates | High   | High   | P2       |
"""

    def _fallback_svg(self, title: str) -> str:
        stages = ["S1\nIdea", "S2\nResearch", "S3\nPRD", "S4\nArch", "S5\nBackend",
                  "S6\nAPI", "S7\nFrontend", "S8\nTests", "S9\nDocker",
                  "S10\nCI/CD", "S11\nDocs", "S12\nMonitor", "S13\nDesign"]
        boxes = ""
        for i, s in enumerate(stages):
            x = 20 + i * 90
            label = s.replace("\n", "&#10;")
            boxes += (f'<rect x="{x}" y="60" width="80" height="50" rx="8" '
                      f'fill="rgba(99,102,241,0.3)" stroke="#00f2ff" stroke-width="1.5"/>'
                      f'<text x="{x+40}" y="82" text-anchor="middle" fill="#00f2ff" '
                      f'font-size="10" font-family="Inter,sans-serif">{s.split(chr(10))[0]}</text>'
                      f'<text x="{x+40}" y="98" text-anchor="middle" fill="#e2e8f0" '
                      f'font-size="9" font-family="Inter,sans-serif">{s.split(chr(10))[1] if chr(10) in s else ""}</text>')
            if i < len(stages) - 1:
                boxes += f'<line x1="{x+80}" y1="85" x2="{x+90}" y2="85" stroke="#00f2ff" stroke-width="1.5" marker-end="url(#arrow)"/>'

        return f"""<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1200 170" width="1200" height="170">
  <defs>
    <marker id="arrow" markerWidth="6" markerHeight="6" refX="3" refY="3" orient="auto">
      <path d="M0,0 L6,3 L0,6 Z" fill="#00f2ff"/>
    </marker>
  </defs>
  <rect width="1200" height="170" fill="#050a0f"/>
  <text x="600" y="30" text-anchor="middle" fill="#00f2ff" font-size="14"
    font-family="Inter,sans-serif" font-weight="bold">{title} — Beta Swarm Pipeline</text>
  {boxes}
</svg>"""
