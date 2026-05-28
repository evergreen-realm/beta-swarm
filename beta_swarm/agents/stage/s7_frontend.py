"""
S7 Frontend Agent — unified with HuashuSkill design generation.
s7_frontend.py and s7_frontend_huashu.py are merged here.
S7FrontendHuashuAgent is an alias for backward compat.
"""
import json, re, os, logging
from typing import Dict, Any, List
from beta_swarm.agents.base import BaseAgent

logger = logging.getLogger(__name__)


class S7FrontendAgent(BaseAgent):
    """
    Stage 7 — Frontend Generation.
    Component-by-component React generation via LLM (GitHub Copilot hint),
    plus Huashu design assets (prototype HTML, PPTX spec, SVG infographic).
    """

    def __init__(self, brain=None):
        super().__init__("s7_frontend", "Frontend Agent", "Stage 7: Frontend Generation", brain)
        self._huashu = None

    def _get_huashu(self):
        if self._huashu is None:
            try:
                from beta_swarm.tools.huashu.huashu_skill import HuashuSkill
                self._huashu = HuashuSkill()
            except Exception:
                pass
        return self._huashu

    def _get_default_next_stage(self):
        return "s8_testing"

    # ------------------------------------------------------------------ #
    def execute(self, task: Dict[str, Any]) -> Dict[str, Any]:
        project_id   = task.get("project_id", "default")
        project_path = task.get("project_path", f"./projects/{project_id}")

        s3_out      = task.get("s3_prd", {})
        prd         = s3_out.get("prd") or task.get("prd") or {}
        s4_out      = task.get("s4_architecture", {})
        architecture = s4_out.get("architecture") or task.get("architecture", {})
        s6_out      = task.get("s6_api", {})
        api_client_code = s6_out.get("api_client_code") or task.get("api_client_code", "")
        api_contracts   = s6_out.get("api_contracts") or architecture.get("api_contracts", [])

        self._log_handover(f"S7 started. Project={project_id}")

        ui_req   = prd.get("ui_ux_requirements", {})
        design_system  = ui_req.get("design_system", "dark mode glassmorphism")
        components_needed = ui_req.get("components", ["App", "Login", "Dashboard", "ItemList"])
        title    = prd.get("metadata", {}).get("title", "App")
        tech     = prd.get("tech_stack_recommendation", {})
        is_react = isinstance(tech, dict) and "react" in str(tech.get("frontend", "")).lower()

        frontend_dir = os.path.join(project_path, "frontend")
        src_dir      = os.path.join(frontend_dir, "src")
        comp_dir     = os.path.join(src_dir, "components")
        os.makedirs(comp_dir, exist_ok=True)

        # 1. Write APIClient.js (from S6 or fallback)
        if api_client_code:
            with open(os.path.join(src_dir, "APIClient.js"), "w", encoding="utf-8") as f:
                f.write(api_client_code)

        generated_files = []

        # 2. Generate components one-by-one (GitHub Copilot hint via router)
        for component in components_needed:
            prompt = f"""You are an expert frontend developer (GitHub Copilot mode).
Generate a React component named {component} for "{title}".
Design system: {design_system}
Use APIClient from '../APIClient.js' for data fetching.
Output ONLY the JSX/JavaScript component code — no markdown, no explanations."""

            component_code = self._call_llm(prompt, task_type="frontend", model_hint="copilot")

            # Strip markdown wrapper if LLM returns one
            m = re.search(r'```(?:jsx?|javascript|tsx?)?\s*\n?(.*?)\n?```', component_code, re.DOTALL)
            if m:
                component_code = m.group(1).strip()

            if len(component_code) < 50:
                component_code = self._fallback_component(component, title, design_system)

            file_path = os.path.join(comp_dir, f"{component}.jsx")
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(component_code)
            generated_files.append(f"frontend/src/components/{component}.jsx")
            self._log_handover(f"S7 generated component: {component}")

        # 3. App.js
        app_prompt = f"""Generate a complete React App.js that imports and renders these components:
{components_needed}
Use APIClient for data. Design: {design_system}. Output ONLY code."""
        app_code = self._call_llm(app_prompt, task_type="frontend")
        m = re.search(r'```(?:jsx?|javascript)?\s*\n?(.*?)\n?```', app_code, re.DOTALL)
        if m:
            app_code = m.group(1).strip()
        if len(app_code) < 50:
            app_code = self._fallback_app(components_needed, title)
        with open(os.path.join(src_dir, "App.js"), "w", encoding="utf-8") as f:
            f.write(app_code)
        generated_files.append("frontend/src/App.js")

        # 4. index.html
        html_prompt = "Generate a minimal index.html that mounts a React app to div#root. Output ONLY HTML code."
        html_code = self._call_llm(html_prompt, task_type="frontend")
        m = re.search(r'```(?:html)?\s*\n?(.*?)\n?```', html_code, re.DOTALL)
        if m:
            html_code = m.group(1).strip()
        if len(html_code) < 50:
            html_code = self._fallback_html(title)
        with open(os.path.join(frontend_dir, "index.html"), "w", encoding="utf-8") as f:
            f.write(html_code)
        generated_files.append("frontend/index.html")

        # 5. styles.css
        css_prompt = f"Generate modern CSS for a '{design_system}' theme. Include glassmorphism cards, dark background, smooth transitions. Output ONLY CSS."
        css_code = self._call_llm(css_prompt, task_type="frontend")
        m = re.search(r'```(?:css)?\s*\n?(.*?)\n?```', css_code, re.DOTALL)
        if m:
            css_code = m.group(1).strip()
        if len(css_code) < 30:
            css_code = self._fallback_css()
        with open(os.path.join(src_dir, "styles.css"), "w", encoding="utf-8") as f:
            f.write(css_code)
        generated_files.append("frontend/src/styles.css")

        # 6. Huashu design assets (prototype + infographic)
        huashu_files = self._run_huashu(project_path, title, prd, architecture)
        generated_files.extend(huashu_files)

        # 7. Frontend Sentry (light static gate)
        self._run_frontend_sentry(frontend_dir, api_contracts)

        preview_url = f"http://localhost:5000/preview/{project_id}"
        self._log_handover(f"S7 completed. {len(generated_files)} files. Preview: {preview_url}")

        os.makedirs(f"./projects/{project_id}", exist_ok=True)
        artifact_path = f"./projects/{project_id}/s7_frontend_output.json"
        artifact = {"files": generated_files, "preview_url": preview_url, "frontend_dir": frontend_dir}
        with open(artifact_path, "w", encoding="utf-8") as f:
            json.dump(artifact, f, indent=2)

        if self.brain:
            try:
                self.brain.store_fact(self.agent_id, f"Frontend: {len(generated_files)} files", "frontend")
            except Exception:
                pass

        return {
            "status": "complete",
            "frontend_files": generated_files,
            "preview_url": preview_url,
            "artifact": artifact,
            "artifact_path": artifact_path,
            "next_stage": task.get("next_stage") or self._get_default_next_stage()
        }

    # ------------------------------------------------------------------ #
    # Huashu design assets
    # ------------------------------------------------------------------ #
    def _run_huashu(self, project_path: str, title: str, prd: dict, arch: dict) -> List[str]:
        files = []
        huashu = self._get_huashu()
        asset_type = "prototype"

        prompt = f"""Generate a premium single-page dark-mode dashboard HTML/CSS/JS prototype for '{title}'.
Use glassmorphism cards, neon cyan accents on dark background (#050a0f), Inter font.
PRD: {str(prd)[:500]}
Return ONE self-contained HTML file:

[FILE: design/prototype.html]
```html
<!DOCTYPE html>
...full code...
```"""

        try:
            llm_resp = self._call_llm(prompt, task_type="frontend")
            parsed = self._parse_files_from_llm(llm_resp)
            for rel, content in parsed.items():
                abs_path = os.path.join(project_path, rel)
                os.makedirs(os.path.dirname(abs_path), exist_ok=True)
                with open(abs_path, "w", encoding="utf-8") as f:
                    f.write(content)
                if huashu:
                    huashu.register_asset(rel, content, "html")
                files.append(rel)
        except Exception as e:
            logger.warning(f"[S7] Huashu prototype generation failed: {e}")
            # Write a guaranteed fallback prototype
            fallback = self._fallback_prototype(title)
            rel = "design/prototype.html"
            abs_path = os.path.join(project_path, rel)
            os.makedirs(os.path.dirname(abs_path), exist_ok=True)
            with open(abs_path, "w", encoding="utf-8") as f:
                f.write(fallback)
            if huashu:
                huashu.register_asset(rel, fallback, "html")
            files.append(rel)

        # PPTX spec
        try:
            pptx_path = os.path.join(project_path, "design", "deck.pptx")
            os.makedirs(os.path.dirname(pptx_path), exist_ok=True)
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            slide.placeholders[1].text = "Beta Swarm Autonomous Design Stack"
            prs.save(pptx_path)
            if huashu:
                huashu.export_pptx("deck", pptx_path)
            files.append("design/deck.pptx")
        except Exception:
            # Markdown fallback
            md_path = os.path.join(project_path, "design", "deck_spec.md")
            os.makedirs(os.path.dirname(md_path), exist_ok=True)
            with open(md_path, "w", encoding="utf-8") as f:
                f.write(f"# {title}\n\n## Slide 1: Title\n\n## Slide 2: Problem\n\n## Slide 3: Architecture\n")
            files.append("design/deck_spec.md")

        return files

    # ------------------------------------------------------------------ #
    # Frontend Sentry (light gate)
    # ------------------------------------------------------------------ #
    def _run_frontend_sentry(self, frontend_dir: str, api_contracts: list):
        required = ["index.html"]
        missing = [f for f in required if not os.path.exists(os.path.join(frontend_dir, f))]
        if missing:
            logger.warning(f"[S7 Sentry] Missing files: {missing}")
        else:
            logger.info("[S7 Sentry] Light gate passed — all required files present")

    # ------------------------------------------------------------------ #
    # Fallback generators
    # ------------------------------------------------------------------ #
    def _fallback_component(self, name: str, title: str, design: str) -> str:
        return f"""import React, {{ useState, useEffect }} from 'react';

const {name} = () => {{
  const [data, setData] = useState([]);

  useEffect(() => {{
    fetch('/api/v1/items/')
      .then(r => r.ok ? r.json() : [])
      .then(d => setData(Array.isArray(d) ? d : []))
      .catch(() => setData([]));
  }}, []);

  return (
    <div className="component-{name.lower()}" style={{{{
      background: 'rgba(255,255,255,0.05)',
      backdropFilter: 'blur(10px)',
      borderRadius: '12px',
      padding: '1.5rem',
      border: '1px solid rgba(255,255,255,0.1)'
    }}}}>
      <h2>{name}</h2>
      <ul>
        {{data.map((item, i) => <li key={{item.id || i}}>{{item.title || JSON.stringify(item)}}</li>)}}
      </ul>
    </div>
  );
}};

export default {name};
"""

    def _fallback_app(self, components: List[str], title: str) -> str:
        imports = "\n".join([f"import {c} from './components/{c}';" for c in components])
        renders = "\n      ".join([f"<{c} />" for c in components])
        return f"""import React from 'react';
import './styles.css';
{imports}

function App() {{
  return (
    <div className="App">
      <header className="app-header">
        <h1>{title}</h1>
      </header>
      <main>
      {renders}
      </main>
    </div>
  );
}}

export default App;
"""

    def _fallback_html(self, title: str) -> str:
        return f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1.0" />
  <title>{title}</title>
</head>
<body>
  <div id="root"></div>
  <script src="src/App.js" type="module"></script>
</body>
</html>"""

    def _fallback_css(self) -> str:
        return """*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
body { background: #050a0f; color: #e2e8f0; font-family: 'Inter', sans-serif; min-height: 100vh; }
.App { display: flex; flex-direction: column; min-height: 100vh; }
.app-header { background: rgba(255,255,255,0.05); backdrop-filter: blur(10px);
  padding: 1rem 2rem; border-bottom: 1px solid rgba(255,255,255,0.1); }
.app-header h1 { font-size: 1.5rem; color: #00f2ff; }
main { padding: 2rem; display: grid; gap: 1.5rem; }
.component-dashboard, .component-itemlist, .component-login {
  background: rgba(255,255,255,0.05); backdrop-filter: blur(10px);
  border-radius: 12px; padding: 1.5rem; border: 1px solid rgba(255,255,255,0.1);
  transition: border-color 0.2s ease; }
.component-dashboard:hover, .component-itemlist:hover { border-color: rgba(0,242,255,0.3); }
button { background: #00f2ff; color: #050a0f; border: none; padding: 0.5rem 1.2rem;
  border-radius: 6px; cursor: pointer; font-weight: 600; transition: opacity 0.2s; }
button:hover { opacity: 0.85; }
"""

    def _fallback_prototype(self, title: str) -> str:
        return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"/>
<title>{title} — Prototype</title>
<style>
  body{{background:#050a0f;color:#e2e8f0;font-family:Inter,sans-serif;display:flex;
       align-items:center;justify-content:center;min-height:100vh;margin:0}}
  .card{{background:rgba(255,255,255,.05);backdrop-filter:blur(12px);border-radius:16px;
         padding:2.5rem;max-width:420px;width:100%;border:1px solid rgba(0,242,255,.2)}}
  h1{{color:#00f2ff;margin-bottom:1rem}}
  .btn{{background:#00f2ff;color:#050a0f;border:none;padding:.6rem 1.4rem;
        border-radius:8px;cursor:pointer;font-weight:700;margin-top:1rem}}
</style>
</head>
<body>
<div class="card">
  <h1>{title}</h1>
  <p>Beta Swarm autonomous prototype.</p>
  <button class="btn" onclick="fetch('/api/v1/items/').then(r=>r.json()).then(d=>alert(JSON.stringify(d)))">
    Fetch Data
  </button>
</div>
</body>
</html>"""

    def _parse_files_from_llm(self, text: str) -> dict:
        files = {}
        patterns = [
            r"\[FILE:\s*([^\]]+)\]\s*```[a-zA-Z0-9]*\n?(.*?)\n?```",
            r"###\s*([a-zA-Z0-9_./-]+\.[a-zA-Z0-9]+)\s*\n```[a-zA-Z0-9]*\n?(.*?)\n?```",
        ]
        for pat in patterns:
            for m in re.finditer(pat, text, re.DOTALL | re.IGNORECASE):
                path = m.group(1).strip().replace("..", "").strip("/\\")
                content = m.group(2).strip()
                if path and content and len(path) < 100:
                    files[path] = content
        return files


# Backward-compat alias
S7FrontendHuashuAgent = S7FrontendAgent
