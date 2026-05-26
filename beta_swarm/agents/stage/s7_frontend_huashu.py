# s7_frontend_huashu.py - Stage 7 Huashu Design & Frontend Generation Agent
import os
import sys
import logging
from typing import Dict, Any
from beta_swarm.agents.base import BaseAgent
from beta_swarm.tools.huashu.huashu_skill import HuashuSkill

logger = logging.getLogger("beta_swarm.s7_frontend_huashu")

class S7FrontendHuashuAgent(BaseAgent):
    def __init__(self, brain=None):
        super().__init__("s7_frontend_huashu", "Huashu Frontend Agent", "Stage 7: Frontend Generation", brain)
        self.huashu = HuashuSkill()

    def execute(self, task: Dict[str, Any]) -> Dict[str, Any]:
        project_path = task.get("project_path") or os.getcwd()
        asset_type = task.get("asset_type") or "prototype"
        project_title = task.get("title") or "Beta Swarm App"
        
        logger.info(f"[{self.name}] Initiating Stage 7 Huashu generation. Asset type: {asset_type}")
        
        # 1. Gather context
        arch = task.get("architecture") or {}
        prd = task.get("prd") or {}
        
        # 2. Process based on asset_type
        generated_files = []
        output_path = ""
        
        if asset_type == "prototype":
            # Generate high-fidelity HTML/CSS/JS prototype
            prompt = f"""
            Generate a high-fidelity, premium, single-page dashboard HTML/CSS/JS prototype for '{project_title}'.
            Use a dark cyberpunk HUD design theme with deep indigo and neon cyan accents.
            Use elegant, harmonious layout with modern typography (Inter).
            
            PRD Requirements: {prd}
            Architecture Contracts: {arch}
            
            Return exactly ONE single self-contained HTML block inside standard markdown brackets like:
            [FILE: frontend/index.html]
            ```html
            ...
            ```
            Do not write stub comments. Provide full responsive grid and hover animations.
            """
            llm_response = self.call_llm([{"role": "user", "content": prompt}])
            parsed_files = self._parse_files(llm_response)
            
            # Export and register
            for rel, content in parsed_files.items():
                abs_path = os.path.join(project_path, rel)
                self._write_file(abs_path, content)
                self.huashu.register_asset(rel, content, "html")
                generated_files.append(rel)
                output_path = abs_path
                
                # Store HTML prototype as artifact so App Preview can render it
                if self.brain:
                    self.brain.store_artifact(
                        agent_id=self.agent_id,
                        project=project_title,
                        stage="S7",
                        data=content
                    )
                
            if not generated_files:
                # Fallback prototype if LLM output fails parsing
                fallback_content = f"<!DOCTYPE html><html><head><title>{project_title}</title></head><body><h1>{project_title} Prototype</h1></body></html>"
                rel = "frontend/index.html"
                abs_path = os.path.join(project_path, rel)
                self._write_file(abs_path, fallback_content)
                self.huashu.register_asset(rel, fallback_content, "html")
                generated_files.append(rel)
                output_path = abs_path
                
                if self.brain:
                    self.brain.store_artifact(
                        agent_id=self.agent_id,
                        project=project_title,
                        stage="S7",
                        data=fallback_content
                    )
                
        elif asset_type == "pptx":
            # Generate slide presentation specifications
            prompt = f"""
            Design a premium slide deck outline for '{project_title}'.
            Provide 5 detailed slide specifications (Title, Problem, Architecture, Swarm Flow, Conclusion)
            including layout grids, colors (indigo/gold/white), and full text copy.
            
            Return exactly ONE slide design markdown document like:
            [FILE: design/deck_spec.md]
            ```markdown
            ...
            ```
            """
            llm_response = self.call_llm([{"role": "user", "content": prompt}])
            parsed_files = self._parse_files(llm_response)
            
            for rel, content in parsed_files.items():
                abs_path = os.path.join(project_path, rel)
                self._write_file(abs_path, content)
                self.huashu.register_asset(rel, content, "pptx")
                generated_files.append(rel)
                output_path = abs_path
                
            # Attempt to use python-pptx for a physical PPTX file
            pptx_path = os.path.join(project_path, "design/deck.pptx")
            os.makedirs(os.path.dirname(pptx_path), exist_ok=True)
            try:
                from pptx import Presentation
                prs = Presentation()
                # Create basic slides based on specifications
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = project_title
                subtitle.text = "Beta Swarm Autonomous Design Stack"
                prs.save(pptx_path)
                self.huashu.export_pptx("deck", pptx_path)
                generated_files.append("design/deck.pptx")
                logger.info("Successfully generated physical deck.pptx using python-pptx.")
            except ImportError:
                # Safe fallback outline if pptx not installed
                logger.warning("python-pptx not available. Generated markdown slide specifications deck_spec.md instead.")
                
        elif asset_type == "infographic":
            # Generate premium raw SVG infographic mapping the swarm workflow
            prompt = f"""
            Generate a complete, raw, beautiful SVG infographic representing the '{project_title}' Swarm Pipeline.
            Include 8 layers (L0 to L8) with clean pulsing gradient lines and modern text badges.
            Use a dark dark background #050a0f, glow cyan lines (#00f2ff), and custom SVG styling.
            
            Return exactly ONE raw SVG block inside brackets like:
            [FILE: design/infographic.svg]
            ```xml
            <svg ...>
            ...
            </svg>
            ```
            """
            llm_response = self.call_llm([{"role": "user", "content": prompt}])
            parsed_files = self._parse_files(llm_response)
            
            for rel, content in parsed_files.items():
                abs_path = os.path.join(project_path, rel)
                self._write_file(abs_path, content)
                self.huashu.register_asset(rel, content, "svg")
                generated_files.append(rel)
                output_path = abs_path
                
            if not generated_files:
                # Fallback clean SVG
                fallback_svg = f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"><rect width="100" height="100" fill="#050a0f"/><text x="10" y="50" fill="#00f2ff">{project_title}</text></svg>'
                rel = "design/infographic.svg"
                abs_path = os.path.join(project_path, rel)
                self._write_file(abs_path, fallback_svg)
                self.huashu.register_asset(rel, fallback_svg, "svg")
                generated_files.append(rel)
                output_path = abs_path
                
        # 3. Store result in Brain if present
        if self.brain:
            self.brain.store_agent_memory(
                self.agent_id,
                f"Generated premium frontend asset {asset_type} for project {project_title}.",
                "design_asset"
            )
            
        logger.info(f"[{self.name}] Stage 7 Huashu generation complete. Generated: {generated_files}")
        
        return {
            "status": "complete",
            "output_path": output_path,
            "asset_type": asset_type,
            "generated_files": generated_files,
            "next_stage": "s8_testing"
        }
